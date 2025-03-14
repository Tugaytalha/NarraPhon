# TODO List:
## 1. Optimize imports
## 3. Delete unnecessary libs from env
## 5. Test multitreading
## 10. Not to do list while prepeare input
## 4. Add subtitle sync from mp3 + txt to srt

import gc
import glob
import re
import shutil
from slide_main import create_presentation as slide_generator

import gradio as gr
import scipy.io.wavfile as wavfile
from Models import *
from Utils import *
from models import *
from moviepy.editor import *
from moviepy.video.tools.subtitles import SubtitlesClip
from utils import *
import sox
import librosa
import soundfile as sf
from shutil import copyfile

torch.manual_seed(0)
torch.backends.cudnn.benchmark = False
torch.backends.cudnn.deterministic = True

import random

random.seed(0)

import numpy as np

np.random.seed(0)

# load packages
import yaml
import numpy as np
import torch
import torchaudio
import librosa
from nltk.tokenize import word_tokenize
import nltk

nltk.download('punkt')
from models import *
from utils import *
from text_utils import TextCleaner
from txtsplit import txtsplit
from pydub import AudioSegment
from pptx import Presentation
import zipfile
import os
from num2words import num2words
import pysrt

textclenaer = TextCleaner()

to_mel = torchaudio.transforms.MelSpectrogram(
    n_mels=80, n_fft=2048, win_length=1200, hop_length=300)
mean, std = -4, 4

device = 'cuda' if torch.cuda.is_available() else 'cpu'


def length_to_mask(lengths):
    mask = torch.arange(lengths.max()).unsqueeze(0).expand(lengths.shape[0], -1).type_as(lengths)
    mask = torch.gt(mask + 1, lengths.unsqueeze(1))
    return mask


def preprocess(wave):
    wave_tensor = torch.from_numpy(wave).float()
    mel_tensor = to_mel(wave_tensor)
    mel_tensor = (torch.log(1e-5 + mel_tensor.unsqueeze(0)) - mean) / std
    return mel_tensor


def compute_style(path):
    wave, sr = librosa.load(path, sr=24000)
    audio, index = librosa.effects.trim(wave, top_db=30)
    if sr != 24000:
        audio = librosa.resample(audio, sr, 24000)
    mel_tensor = preprocess(audio).to(device)

    with torch.no_grad():
        ref_s = model.style_encoder(mel_tensor.unsqueeze(1))
        ref_p = model.predictor_encoder(mel_tensor.unsqueeze(1))

    return torch.cat([ref_s, ref_p], dim=1)


# load phonemizer
import phonemizer

global_phonemizer = phonemizer.backend.EspeakBackend(language='en-us', preserve_punctuation=True, with_stress=True)

config = yaml.safe_load(open("Models/LibriTTS/config.yml"))

# load pretrained ASR model
ASR_config = config.get('ASR_config', False)
ASR_path = config.get('ASR_path', False)
text_aligner = load_ASR_models(ASR_path, ASR_config)

# load pretrained F0 model
F0_path = config.get('F0_path', False)
pitch_extractor = load_F0_models(F0_path)

# load BERT model
from Utils.PLBERT.util import load_plbert

BERT_path = config.get('PLBERT_dir', False)
plbert = load_plbert(BERT_path)

model_params = recursive_munch(config['model_params'])
model = build_model(model_params, text_aligner, pitch_extractor, plbert)
_ = [model[key].eval() for key in model]
_ = [model[key].to(device) for key in model]

params_whole = torch.load("Models/LibriTTS/epochs_2nd_00020.pth", map_location='cpu')
params = params_whole['net']

for key in model:
    if key in params:
        print('%s loaded' % key)
        try:
            model[key].load_state_dict(params[key])
        except:
            from collections import OrderedDict

            state_dict = params[key]
            new_state_dict = OrderedDict()
            for k, v in state_dict.items():
                name = k[7:]  # remove `module.`
                new_state_dict[name] = v
            # load params
            model[key].load_state_dict(new_state_dict, strict=False)
#             except:
#                 _load(params[key], model[key])
_ = [model[key].eval() for key in model]

from Modules.diffusion.sampler import DiffusionSampler, ADPM2Sampler, KarrasSchedule

sampler = DiffusionSampler(
    model.diffusion.diffusion,
    sampler=ADPM2Sampler(),
    sigma_schedule=KarrasSchedule(sigma_min=0.0001, sigma_max=3.0, rho=9.0),  # empirical parameters
    clamp=False
)


def inference(text, ref_s, alpha=0.3, beta=0.7, diffusion_steps=5, embedding_scale=1):
    text = text.strip()
    ps = global_phonemizer.phonemize([text])
    ps = word_tokenize(ps[0])
    ps = ' '.join(ps)
    tokens = textclenaer(ps)
    tokens.insert(0, 0)
    tokens = torch.LongTensor(tokens).to(device).unsqueeze(0)

    with torch.no_grad():
        input_lengths = torch.LongTensor([tokens.shape[-1]]).to(device)
        text_mask = length_to_mask(input_lengths).to(device)

        t_en = model.text_encoder(tokens, input_lengths, text_mask)
        bert_dur = model.bert(tokens, attention_mask=(~text_mask).int())
        d_en = model.bert_encoder(bert_dur).transpose(-1, -2)

        s_pred = sampler(noise=torch.randn((1, 256)).unsqueeze(1).to(device),
                         embedding=bert_dur,
                         embedding_scale=embedding_scale,
                         features=ref_s,  # reference from the same speaker as the embedding
                         num_steps=diffusion_steps).squeeze(1)

        s = s_pred[:, 128:]
        ref = s_pred[:, :128]

        ref = alpha * ref + (1 - alpha) * ref_s[:, :128]
        s = beta * s + (1 - beta) * ref_s[:, 128:]

        d = model.predictor.text_encoder(d_en,
                                         s, input_lengths, text_mask)

        x, _ = model.predictor.lstm(d)
        duration = model.predictor.duration_proj(x)

        duration = torch.sigmoid(duration).sum(axis=-1)
        pred_dur = torch.round(duration.squeeze()).clamp(min=1)

        pred_aln_trg = torch.zeros(input_lengths, int(pred_dur.sum().data))
        c_frame = 0
        for i in range(pred_aln_trg.size(0)):
            pred_aln_trg[i, c_frame:c_frame + int(pred_dur[i].data)] = 1
            c_frame += int(pred_dur[i].data)

        # encode prosody
        en = (d.transpose(-1, -2) @ pred_aln_trg.unsqueeze(0).to(device))
        if model_params.decoder.type == "hifigan":
            asr_new = torch.zeros_like(en)
            asr_new[:, :, 0] = en[:, :, 0]
            asr_new[:, :, 1:] = en[:, :, 0:-1]
            en = asr_new

        F0_pred, N_pred = model.predictor.F0Ntrain(en, s)

        asr = (t_en @ pred_aln_trg.unsqueeze(0).to(device))
        if model_params.decoder.type == "hifigan":
            asr_new = torch.zeros_like(asr)
            asr_new[:, :, 0] = asr[:, :, 0]
            asr_new[:, :, 1:] = asr[:, :, 0:-1]
            asr = asr_new

        out = model.decoder(asr,
                            F0_pred, N_pred, ref.squeeze().unsqueeze(0))

    return out.squeeze().cpu().numpy()[..., :-50]  # weird pulse at the end of the model, need to be fixed later


def LFinference(text, s_prev, ref_s, alpha=0.3, beta=0.7, t=0.7, diffusion_steps=5, embedding_scale=1):
    text = text.strip()
    ps = global_phonemizer.phonemize([text])
    ps = word_tokenize(ps[0])
    ps = ' '.join(ps)
    ps = ps.replace('``', '"')
    ps = ps.replace("''", '"')

    tokens = textclenaer(ps)
    tokens.insert(0, 0)
    tokens = torch.LongTensor(tokens).to(device).unsqueeze(0)

    with torch.no_grad():
        input_lengths = torch.LongTensor([tokens.shape[-1]]).to(device)
        text_mask = length_to_mask(input_lengths).to(device)

        t_en = model.text_encoder(tokens, input_lengths, text_mask)
        bert_dur = model.bert(tokens, attention_mask=(~text_mask).int())
        d_en = model.bert_encoder(bert_dur).transpose(-1, -2)

        s_pred = sampler(noise=torch.randn((1, 256)).unsqueeze(1).to(device),
                         embedding=bert_dur,
                         embedding_scale=embedding_scale,
                         features=ref_s,  # reference from the same speaker as the embedding
                         num_steps=diffusion_steps).squeeze(1)

        if s_prev is not None:
            # convex combination of previous and current style
            s_pred = t * s_prev + (1 - t) * s_pred

        s = s_pred[:, 128:]
        ref = s_pred[:, :128]

        ref = alpha * ref + (1 - alpha) * ref_s[:, :128]
        s = beta * s + (1 - beta) * ref_s[:, 128:]

        s_pred = torch.cat([ref, s], dim=-1)

        d = model.predictor.text_encoder(d_en,
                                         s, input_lengths, text_mask)

        x, _ = model.predictor.lstm(d)
        duration = model.predictor.duration_proj(x)

        duration = torch.sigmoid(duration).sum(axis=-1)
        pred_dur = torch.round(duration.squeeze()).clamp(min=1)

        pred_aln_trg = torch.zeros(input_lengths, int(pred_dur.sum().data))
        c_frame = 0
        for i in range(pred_aln_trg.size(0)):
            pred_aln_trg[i, c_frame:c_frame + int(pred_dur[i].data)] = 1
            c_frame += int(pred_dur[i].data)

        # encode prosody
        en = (d.transpose(-1, -2) @ pred_aln_trg.unsqueeze(0).to(device))
        if model_params.decoder.type == "hifigan":
            asr_new = torch.zeros_like(en)
            asr_new[:, :, 0] = en[:, :, 0]
            asr_new[:, :, 1:] = en[:, :, 0:-1]
            en = asr_new

        F0_pred, N_pred = model.predictor.F0Ntrain(en, s)

        asr = (t_en @ pred_aln_trg.unsqueeze(0).to(device))
        if model_params.decoder.type == "hifigan":
            asr_new = torch.zeros_like(asr)
            asr_new[:, :, 0] = asr[:, :, 0]
            asr_new[:, :, 1:] = asr[:, :, 0:-1]
            asr = asr_new

        out = model.decoder(asr,
                            F0_pred, N_pred, ref.squeeze().unsqueeze(0))

    return out.squeeze().cpu().numpy()[...,
           :-100], s_pred  # weird pulse at the end of the model, need to be fixed later


def STinference(text, ref_s, ref_text, alpha=0.3, beta=0.7, diffusion_steps=5, embedding_scale=1):
    text = text.strip()
    ps = global_phonemizer.phonemize([text])
    ps = word_tokenize(ps[0])
    ps = ' '.join(ps)

    tokens = textclenaer(ps)
    tokens.insert(0, 0)
    tokens = torch.LongTensor(tokens).to(device).unsqueeze(0)

    ref_text = ref_text.strip()
    ps = global_phonemizer.phonemize([ref_text])
    ps = word_tokenize(ps[0])
    ps = ' '.join(ps)

    ref_tokens = textclenaer(ps)
    ref_tokens.insert(0, 0)
    ref_tokens = torch.LongTensor(ref_tokens).to(device).unsqueeze(0)

    with torch.no_grad():
        input_lengths = torch.LongTensor([tokens.shape[-1]]).to(device)
        text_mask = length_to_mask(input_lengths).to(device)

        t_en = model.text_encoder(tokens, input_lengths, text_mask)
        bert_dur = model.bert(tokens, attention_mask=(~text_mask).int())
        d_en = model.bert_encoder(bert_dur).transpose(-1, -2)

        ref_input_lengths = torch.LongTensor([ref_tokens.shape[-1]]).to(device)
        ref_text_mask = length_to_mask(ref_input_lengths).to(device)
        ref_bert_dur = model.bert(ref_tokens, attention_mask=(~ref_text_mask).int())
        s_pred = sampler(noise=torch.randn((1, 256)).unsqueeze(1).to(device),
                         embedding=bert_dur,
                         embedding_scale=embedding_scale,
                         features=ref_s,  # reference from the same speaker as the embedding
                         num_steps=diffusion_steps).squeeze(1)

        s = s_pred[:, 128:]
        ref = s_pred[:, :128]

        ref = alpha * ref + (1 - alpha) * ref_s[:, :128]
        s = beta * s + (1 - beta) * ref_s[:, 128:]

        d = model.predictor.text_encoder(d_en,
                                         s, input_lengths, text_mask)

        x, _ = model.predictor.lstm(d)
        duration = model.predictor.duration_proj(x)

        duration = torch.sigmoid(duration).sum(axis=-1)
        pred_dur = torch.round(duration.squeeze()).clamp(min=1)

        pred_aln_trg = torch.zeros(input_lengths, int(pred_dur.sum().data))
        c_frame = 0
        for i in range(pred_aln_trg.size(0)):
            pred_aln_trg[i, c_frame:c_frame + int(pred_dur[i].data)] = 1
            c_frame += int(pred_dur[i].data)

        # encode prosody
        en = (d.transpose(-1, -2) @ pred_aln_trg.unsqueeze(0).to(device))
        if model_params.decoder.type == "hifigan":
            asr_new = torch.zeros_like(en)
            asr_new[:, :, 0] = en[:, :, 0]
            asr_new[:, :, 1:] = en[:, :, 0:-1]
            en = asr_new

        F0_pred, N_pred = model.predictor.F0Ntrain(en, s)

        asr = (t_en @ pred_aln_trg.unsqueeze(0).to(device))
        if model_params.decoder.type == "hifigan":
            asr_new = torch.zeros_like(asr)
            asr_new[:, :, 0] = asr[:, :, 0]
            asr_new[:, :, 1:] = asr[:, :, 0:-1]
            asr = asr_new

        out = model.decoder(asr,
                            F0_pred, N_pred, ref.squeeze().unsqueeze(0))

    return out.squeeze().cpu().numpy()[..., :-50]  # weird pulse at the end of the model, need to be fixed later


def split_text(text, max_length=50):
    words = text.split()
    chunks = []
    current_chunk = []

    for i in range(words):
        word = words[i]
        if len(current_chunk) + len(word) + 1 <= max_length:
            current_chunk.append(word)
        else:
            current_chunk.append("for")
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks


def convert_numbers_to_words(text):
    # Regular expression to find numbers in the text
    number_pattern = re.compile(r'\b\d+\b')

    # Function to replace each number with its word equivalent
    def replace_number(match):
        number = int(match.group(0))
        return num2words(number)

    # Replace all numbers in the text using the regular expression
    result_text = re.sub(number_pattern, replace_number, text)
    return result_text


def concatenate_srt_files(srt_files, output_file="./extracted/concatenated.srt"):
    """
    Concatenate multiple SRT files into a single SRT file.

    :param srt_files: List of paths to the SRT files to be concatenated.
    :param output_file: Path to the output SRT file.
    """
    final_subtitles = pysrt.SubRipFile()
    current_time_offset = pysrt.SubRipTime(0, 0, 0, 0)
    last_count = 0

    for srt_file in srt_files:
        subs = pysrt.open(srt_file)

        for sub in subs:
            # Update the start and end times of the subtitle
            sub.shift(seconds=current_time_offset.seconds,
                      minutes=current_time_offset.minutes,
                      hours=current_time_offset.hours,
                      milliseconds=current_time_offset.milliseconds)
            # Update the subtitle number
            last_count += 1
            sub.index = last_count
            # Append the subtitle to the final list
            final_subtitles.append(sub)

        # Update the current time offset to the end time of the last subtitle + 1 second
        last_sub = subs[-1]
        current_time_offset = last_sub.end + pysrt.SubRipTime(0, 0, 1, 0)

        # Update the last count
        last_count = final_subtitles[-1].index

    final_subtitles.save(output_file, encoding='utf-8')

def generate_recursively(audio_file, directory, speed, alpha, beta, diffusion_steps, embedding_scale,
                         file_encoding="utf-8"):
    # Use glob to find all .txt files recursively
    txt_files = glob.glob(os.path.join(directory, '**', '*.txt'), recursive=True)
    txt_files.sort(key=natural_keys)

    # Save a concatenated text file
    with open(directory + "/concatenated.txt", "w", encoding=file_encoding) as file:
        for txt_file in txt_files:
            with open(txt_file, "r", encoding=file_encoding) as f:
                file.write(f.read() + "\n")

    # Create generated_voices directory if it doesn't exist
    if not os.path.exists(directory + "/generated_voices"):
        os.makedirs(directory + "/generated_voices")

    # Create generated_subtitle directory if it doesn't exist
    if not os.path.exists(directory + "/generated_subtitle"):
        os.makedirs(directory + "/generated_subtitle")

    print("Generating speech for all text files...")
    for txt_file in txt_files:
        try:
            with open(txt_file, 'r', encoding=file_encoding) as file:
                content = file.read()
                output = generate_speech(audio_file, content, speed, alpha, beta, diffusion_steps, embedding_scale)

                if output:
                    # Get the output file name with convert txt to mp3 and adding generated_voices to the path
                    output_file = txt_file.replace(".txt", ".mp3").replace(directory + "/",
                                                                           directory + "/generated_voices/")

                    # Save the generated speech as an MP3 file
                    wav_file = "temp.wav"
                    wavfile.write(wav_file, *output[0])
                    audio = AudioSegment.from_file(wav_file)
                    audio.export(output_file, format="mp3")

                    # Delete the .wav file
                    os.remove(wav_file)

                    # Delete the .txt file
                    os.remove(txt_file)

                    # Generate subtitle
                    subtitle_directory = directory + "/generated_subtitle"
                    subtitle_file = txt_file.replace(".txt", ".srt").replace(directory + "/", subtitle_directory + "/")

                    # Create subtitle
                    os.system(
                        f"whisper {output_file} --model small --language English --max_line_count=1 --max_line_width=80 --word_timestamps=True --output_dir {subtitle_directory} --output_format=srt")

                    # Fix known subtitle mistakes
                    correct_known_mistakes(subtitle_file)


        except Exception as e:
            print(f"Error reading {txt_file}: {e}")

        gc.collect()

        print("Audio generated")

    # Concatenate the generated audio files
    generated_files = glob.glob(os.path.join(directory + "/generated_voices", '**', '*.mp3'), recursive=True)
    # Sort the generated files
    generated_files.sort(key=natural_keys)
    audio = AudioSegment.from_file(generated_files[0])
    # Calculate the duration of the first audio file
    for file in generated_files[1:]:
        audio += AudioSegment.from_file(file)


    # Read and sort the generated subtitle files
    generated_files = glob.glob(os.path.join(directory + "/generated_subtitle", '**', '*.srt'), recursive=True)
    generated_files.sort(key=natural_keys)

    # Concatenate the generated subtitle files
    concatenate_srt_files(generated_files, directory + "/concatenated.srt")

    # Export the concatenated audio
    audio.export(directory + "/concatenated.mp3", format="mp3")


def gen_from_text(audio_file, text, speed, alpha, beta, diffusion_steps, embedding_scale):
    # Generate speech from the text
    voice, message = generate_speech(audio_file, text, speed, alpha, beta, diffusion_steps, embedding_scale)

    # Save the generated speech as an MP3 file
    wav_file = "temp.wav"
    wavfile.write(wav_file, *voice)
    audio = AudioSegment.from_file(wav_file)
    audio.export("output.mp3", format="mp3")

    # # Create subtitle
    # os.system(
    #     f"whisper output.mp3 --model small --language English --max_line_count=2 --max_line_width=60 --word_timestamps=True --output_dir . --output_format=srt")

    # Zip the generated files
    with zipfile.ZipFile("output.zip", 'w', zipfile.ZIP_DEFLATED) as zipf:
        zipf.write("output.mp3")
        # zipf.write("output.srt")

    return "output.zip", "temp.wav", message


def gen_from_srt(audio_file, srt_file, speed, alpha, beta, diffusion_steps, embedding_scale):
    # Load the SRT file
    subs = pysrt.open(srt_file)

    # Merge subtitles that are continuous
    for i in range(len(subs)-2, -1, -1):
        if i < len(subs) - 1:
            if subs[i].end.ordinal == subs[i + 1].start.ordinal:
                subs[i].text += " " + subs[i + 1].text
                subs[i].end = subs[i + 1].end
                subs.pop(i + 1)


    synthesized_audio_list = []
    for i in range(len(subs)):
        # Generate speech for the subtitle text
        voice, message = generate_speech(audio_file, subs[i].text, speed, alpha, beta, diffusion_steps, embedding_scale)
        synthesized_audio_list.append(voice[1])

        # If this is not the last subtitle, add silence for the duration between the end of the current subtitle and the start of the next subtitle
        if i < len(subs) - 1:
            # Calculate the length of the output voice
            voice_duration = len(voice[1]) / voice[0]
            silence_duration = ((subs[i + 1].start.ordinal - subs[i].start.ordinal) / 1000) - voice_duration
            if silence_duration > 0:
                silence = np.zeros(int(silence_duration * 24000)).astype(np.int16)
                synthesized_audio_list.append(silence)

    # Concatenate the synthesized audio
    synthesized_audio = np.concatenate(synthesized_audio_list, axis=0)

    # Save the generated speech as an MP3 file
    wavfile.write("temp.wav", 24000, synthesized_audio)
    audio = AudioSegment.from_file("temp.wav")
    audio.export("output.mp3", format="mp3")

    # Zip the generated files
    with zipfile.ZipFile("output.zip", 'w', zipfile.ZIP_DEFLATED) as zipf:
        zipf.write("output.mp3")

    return "output.zip", "temp.wav", message


def atoi(text):
    return int(text) if text.isdigit() else text


def natural_keys(text):
    '''
    alist.sort(key=natural_keys) sorts in human order
    http://nedbatchelder.com/blog/200712/human_sorting.html
    (See Toothy's implementation in the comments)
    '''
    return [atoi(c) for c in re.split(r'(\d+)', text)]


def correct_known_mistakes(file_path="extracted/concatenated.srt"):
    # Huawei
    correct_mistaken_words(file_path=file_path)
    # &
    srt_file_path = file_path
    with open(srt_file_path, "r") as file:
        content = file.read()
    content = re.sub(r'(?<=\b\w) and (?=\w\b)', '&', content)

    with open(srt_file_path, "w") as file:
        file.write(content)

    # Slash
    correct_mistaken_words(incorrect_words=["-slash-", "slash"], correct_word="/", file_path=file_path)
    # CodeArts
    correct_mistaken_words(incorrect_words=["Code Arts", "codars"], correct_word="CodeArts", file_path=file_path)
    # 6G
    correct_mistaken_words(incorrect_words=["6 G", "6-G"], correct_word="6G", file_path=file_path)
    # 5G
    correct_mistaken_words(incorrect_words=["5 G", "5-G"], correct_word="5G", file_path=file_path)
    # 4G
    correct_mistaken_words(incorrect_words=["4 G", "4-G"], correct_word="4G", file_path=file_path)
    # 3G
    correct_mistaken_words(incorrect_words=["3 G", "3-G"], correct_word="3G", file_path=file_path)
    # 2G
    correct_mistaken_words(incorrect_words=["2 G", "2-G"], correct_word="2G", file_path=file_path)
    # DevSecOps
    correct_mistaken_words(incorrect_words=["DevSecUps"], correct_word="DevSecOps", file_path=file_path)
    # DevOps
    correct_mistaken_words(incorrect_words=["DevUps"], correct_word="DevOps", file_path=file_path)
    # PerfTest
    correct_mistaken_words(incorrect_words=["Perf Test"], correct_word="PerfTest", file_path=file_path)
    # Retry
    correct_mistaken_words(incorrect_words=["REIT"], correct_word="retry", file_path=file_path)

    print("Mistaken words corrected successfully.")


def parse_generate(audio_file, text_input_type, text_input, text_file, srt_input, zip_file, pptx_inp,
                   speed, alpha, beta, diffusion_steps, embedding_scale):
    # Clear the memory
    gc.collect()

    output_dir = "extracted"

    # Delete the generated files if they exist
    if os.path.exists("output.zip"):
        os.remove("output.zip")
    if os.path.exists("generated.zip"):
        os.remove("generated.zip")

    # Make sure the audio file is a WAV file
    if audio_file is not None:
        audio_file = convert_to_wav(audio_file)

        if text_input_type == "Plain Text":
            result = gen_from_text(audio_file, text_input, speed, alpha, beta, diffusion_steps, embedding_scale)

        if text_input_type == "TXT File":
            with open(text_file, "r") as f:
                text_input = f.read()
            result = gen_from_text(audio_file, text_input, speed, alpha, beta, diffusion_steps, embedding_scale)

        elif text_input_type == "SRT File":
            result = gen_from_srt(audio_file, srt_input, speed, alpha, beta, diffusion_steps, embedding_scale)

        elif text_input_type.startswith("ZIP File"):

            # Extract zip file
            with zipfile.ZipFile(zip_file, 'r') as zip_ref:
                zip_ref.extractall(output_dir)

            # generate speech recursively
            generate_recursively(audio_file, output_dir, speed, alpha, beta, diffusion_steps, embedding_scale)

            print("generated all files")

            print("generating subtitle...")
            # # os.system(f"whisper {output_dir}/concatenated.mp3 --model small --language English --max_line_count=2 --max_line_width=60 --word_timestamps=True --output_dir {output_dir}/generated_subtitle --output_format=srt")
            # os.system(
            #     f"whisper {output_dir}/concatenated.mp3 --model small --language English --max_line_count=1 --max_line_width=70 --word_timestamps=True --output_dir {output_dir}/generated_subtitle --output_format=srt")
            #
            # # Fix known subtitle mistakes starting withHuawei's
            # correct_known_mistakes()

            print("subtitle generated")

            if text_input_type == "ZIP FileP":
                print("Creating video...")
                create_video()
                print("Video created")

            # Zip the generated files
            zip_output = "generated.zip"
            with zipfile.ZipFile(zip_output, 'w', zipfile.ZIP_DEFLATED) as zipf:
                # Walk the directory
                for root, dirs, files in os.walk(output_dir):
                    for file in files:
                        # Create the full path of the file
                        file_path = os.path.join(root, file)
                        # Add file to the zip file, preserving the directory structure
                        zipf.write(file_path, os.path.relpath(file_path, output_dir))

            # Delete the extracted directory recursively without error
            os.system(f"rm -rf {output_dir} || true")

            result = zip_output, None, "Success: Speech generated successfully."


        elif text_input_type == "PowerPoint File":
            output_dir = "extracted_notes"
            zip_file_path = extract_notes(pptx_inp, output_dir, "notes.zip")

            output_dir = "extracted"

            # Create output directory if it doesn't exist
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)

            # Create output/images directory if it doesn't exist
            if not os.path.exists(output_dir + "/images"):
                os.makedirs(output_dir + "/images")
            # Convert the PowerPoint slides to images
            os.system(f"soffice --headless --convert-to pdf {pptx_inp} --outdir {os.path.dirname(pptx_inp)}")
            os.system(
                f"pdftoppm -r 150 -jpeg -jpegopt quality=100 {pptx_inp.replace('.pptx', '.pdf')} {output_dir + '/images/slide'}")

            # Delete the PDF file
            os.remove(pptx_inp.replace('.pptx', '.pdf'))

            result = parse_generate(audio_file, "ZIP FileP", None, None, srt_input, zip_file_path, None, speed, alpha, beta,
                                    diffusion_steps,
                                    embedding_scale)

        elif text_input_type == "Prompt":
            print("Generating speech from prompt...")
            # Create output directory if it doesn't exist
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)

            print("output_dir: ", output_dir, "pptx_inp: ", pptx_inp)
            print("Generating speech...")
            # Move pptx file to extracted directory
            shutil.move(pptx_inp, output_dir)

            print("Generating speech recursively...")
            # Update pptx_inp to the new path
            pptx_inp = os.path.join(output_dir, os.path.basename(pptx_inp))

            print("pptx_inp: ", pptx_inp)
            print("Extracting notes...")
            result = parse_generate(audio_file, "PowerPoint File", None, None, srt_input, zip_file, pptx_inp, speed, alpha, beta, diffusion_steps, embedding_scale)

        # Delete the generated files if they exist
        if os.path.exists("output.mp3"):
            os.remove("output.mp3")
        if os.path.exists("notes.zip"):
            os.remove("notes.zip")
        if os.path.exists("extracted_notes"):
            # delete recursively and force delete
            os.system("rm -rf extracted_notes || true")

        return result

    else:
        return None, None, "Error: Audio file is missing."


def correct_mistaken_words(file_path="extracted/concatenated.srt",
                           incorrect_words=["heal way", "hue-away", "Woway", "wo way", "hueaway", "hueAway", "whoaway",
                                            "huawei", "worldway", "who-away", "world way", "world away",
                                            "Heway", "heaway", "he way", "he away",
                                            "raw away", "raw way", "raw-way", "raw-away", "who away", "who-away"],
                           correct_word="Huawei"):
    # Read the content of the SRT file
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()

    # Create a regular expression pattern to find the incorrect words
    pattern = re.compile('|'.join(re.escape(word) for word in incorrect_words), re.IGNORECASE)

    # Replace the incorrect words with the correct word
    corrected_content = pattern.sub(correct_word, content)

    # Write the corrected content back to the SRT file
    with open(file_path, 'w', encoding='utf-8') as file:
        file.write(corrected_content)


def create_video(folder_path="extracted", output_path="extracted/output_video.mp4", font="Huawei-Sans-Bold",
                 font_size=44):
    # Determine the number of files
    num_slides = len(os.listdir(f"{folder_path}/images/"))
    max_len = len(str(num_slides))

    # List of images and corresponding audio files
    image_files = [f"{folder_path}/images/slide-{str(i).zfill(max_len)}.jpg" for i in range(1, num_slides + 1)]
    audio_files = [f"{folder_path}/generated_voices/slide_{i}.mp3" for i in range(1, num_slides + 1)]
    subtitle_files = [f"{folder_path}/generated_subtitle/slide_{i}.srt" for i in range(1, num_slides + 1)]

    # Check if the audio files exist if not replace with none
    for audio_file in audio_files:
        if not os.path.exists(audio_file):
            audio_files[audio_files.index(audio_file)] = None

    # List to hold the video clips
    video_clips = []

    # Create video clips from images and corresponding audio
    for image, audio, subtitle_file in zip(image_files, audio_files, subtitle_files):
        try:
            # Load image
            image_clip = ImageClip(image)

            if audio is not None:
                # Load audio
                audio_clip = AudioFileClip(audio)

                # Set the duration of the image to the duration of the audio + 1 second
                image_clip = image_clip.set_duration(audio_clip.duration + 1)

                # Set the audio to the image
                image_clip = image_clip.set_audio(audio_clip)

                # # Add subtitles
                # generator = lambda txt: TextClip(txt, font=font, fontsize=font_size, color='white',
                #                                  stroke_color='black',
                #                                  stroke_width=2.8)
                # subs = SubtitlesClip(subtitle_file, generator)
                # subtitles = SubtitlesClip(subs, generator)
                #
                # image_clip = CompositeVideoClip([image_clip, subtitles.set_position(("center", 0.9), relative=True)])

            else:
                # Set the duration of the image to 3 seconds
                image_clip = image_clip.set_duration(3)

            # Append the clip to the list
            video_clips.append(image_clip)

        except Exception as e:
            print(f"Error creating video clip: {e}")

    # Concatenate all video clips
    final_video = concatenate_videoclips(video_clips)

    # Determine thread number
    thread_count = os.cpu_count() - 1 if os.cpu_count() > 1 else 1

    # # Determine if h254_nvenc is available
    # if os.system("ffmpeg -encoders | grep h264_nvenc ") == 0:
    #     codec = "h264_nvenc "
    # # elif os.system("ffmpeg -encoders | grep h264_videotoolbox") == 0:
    # #     codec = "h264_videotoolbox"
    # # else:
    # #     codec = "libx264"

    try:
        # Write the final output
        final_video.write_videofile(output_path, fps=24, codec="libx264", threads=thread_count, audio_codec="aac")
    except Exception as e:
        print(f"Error writing video: {e}")


def generate_speech(audio_file, text_input, speed, alpha, beta, diffusion_steps, embedding_scale):
    try:
        if audio_file is not None and text_input:
            # Making sure it is a wav file
            audio_file = convert_to_wav(audio_file)

            ref_s = compute_style(audio_file)
            text_chunks = txtsplit(text_input)

            synthesized_audio_list = []
            for chunk in text_chunks:
                # Replace Huawei to Whoaway
                chunk = chunk.replace("Huawei", "Whoaway")
                chunk = chunk.replace("huawei", "whoaway")
                chunk = chunk.replace("HUAWEI", "whoaway")
                chunk = chunk.replace("-", " ")
                chunk = convert_numbers_to_words(chunk)
                pFlag = False

                # Check if the chunk ends with a punctuation
                if chunk[-1] not in [".", "!", "?", ",", ";", ":"]:
                    pFlag = True
                    chunk = chunk + "."
                print(chunk)
                synthesized_audio_chunk = inference(chunk, ref_s, alpha, beta, diffusion_steps,
                                                    embedding_scale)  # [:-int(0.40 * 24000)] # Delete last 350 ms which says garbage "for" for fixing weird pulse at the end

                if pFlag:
                    synthesized_audio_chunk = synthesized_audio_chunk[:-int(0.5 * 24000)]

                synthesized_audio_list.append(synthesized_audio_chunk)

            synthesized_audio = np.concatenate(synthesized_audio_list, axis=0)

            # Convert to 16-bit PCM
            synthesized_audio = (synthesized_audio * 32767).astype(np.int16)

            # Write the synthesized audio to a temp file
            wavfile.write("temp.wav", 24000, synthesized_audio)

            # Speed up with sox
            if speed >= 1:
                # Create a Transformer object
                tfm = sox.Transformer()

                # Set the tempo change (speed change)
                tfm.tempo(speed)

                # Apply the transformation and save the output file
                tfm.build("temp.wav", "temp2.wav")

            # Slow down with sox
            elif speed < 1:
                # Create a Transformer object
                tfm = sox.Transformer()

                # Set the tempo change (speed change)
                tfm.tempo(speed)

                # Apply the transformation and save the output file
                tfm.build("temp.wav", "temp2.wav")

            # Just create temp2.wav
            else:
                # Copy the temp file
                copyfile("temp.wav", "temp2.wav")

            # Read back with pydub
            synthesized_audio = AudioSegment.from_file("temp2.wav")

            # Convert to numpy array
            synthesized_audio = np.array(synthesized_audio.get_array_of_samples())

            # Delete the temp file
            os.remove("temp.wav")
            os.remove("temp2.wav")

            return (24000, synthesized_audio), "Success: Speech generated successfully."

        return None, "Error: Audio file or text input is missing."
    except Exception as e:
        print(e)
        return None, str(e)


def convert_to_wav(input_file, output_file=None):
    # Determine the output file name if not provided
    if not output_file:
        base, ext = os.path.splitext(input_file)
        output_file = base + ".wav"

    # Check if the input file is a WAV file
    if not input_file.endswith(".wav"):
        # Load the audio file
        audio = AudioSegment.from_file(input_file)

        # Export as WAV
        audio.export(output_file, format="wav")

    return output_file


def extract_notes(pptx_path, output_dir, zip_file_path):
    # Create the output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Open the PowerPoint file
    presentation = Presentation(pptx_path)

    # List to hold the paths of the created text files
    text_files = []

    # Extract notes and slides from each slide and save to individual text files
    for slide_number, slide in enumerate(presentation.slides, start=1):
        txt_file_path = os.path.join(output_dir, f"slide_{slide_number}.txt")
        with open(txt_file_path, "w", encoding="utf-8") as txt_file:
            # Extract the slide notes
            if slide.notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                txt_file.write(notes_text)
            else:
                txt_file.write("No notes.")

            text_files.append(txt_file_path)

    # Create a zip file and add all the text files to it
    with zipfile.ZipFile(zip_file_path, 'w') as zipf:
        for text_file in text_files:
            zipf.write(text_file, os.path.basename(text_file))

    return zip_file_path


def update_input_fields(input_type):
    if input_type == "Plain Text":
        return (gr.update(visible=True), gr.update(visible=False), gr.update(visible=False),
                gr.update(visible=False), gr.update(visible=False),
                gr.update(visible=True), gr.update(visible=True))
    elif input_type == "TXT File":
        return (gr.update(visible=False), gr.update(visible=True), gr.update(visible=False),
                gr.update(visible=False), gr.update(visible=False),
                gr.update(visible=True), gr.update(visible=True))
    elif input_type == "SRT File":
        return (gr.update(visible=False), gr.update(visible=False), gr.update(visible=True),
                gr.update(visible=False), gr.update(visible=False),
                gr.update(visible=True), gr.update(visible=True))
    elif input_type == "ZIP File":
        return (gr.update(visible=False), gr.update(visible=False), gr.update(visible=False),
                gr.update(visible=True), gr.update(visible=False),
                gr.update(visible=False), gr.update(visible=True))
    elif input_type == "PowerPoint File":
        return (gr.update(visible=False), gr.update(visible=False), gr.update(visible=False),
                gr.update(visible=False), gr.update(visible=True),
                gr.update(visible=False), gr.update(visible=True))


# iface = gr.Interface(
#    fn=generate_speech,
#    inputs=[
#        gr.Audio(label="Upload a reference audio file", type="filepath"),
#        txt_radio,
#        txt_box,
#        txt_inp,
#        zip_inp,
#        gr.Slider(minimum=0.0, maximum=1.0, value=0.3, label="Alpha"),
#        gr.Slider(minimum=0.0, maximum=1.0, value=0.7, label="Beta"),
#        gr.Slider(minimum=1, maximum=10, value=5, label="Diffusion Steps"),
#        gr.Slider(minimum=0.0, maximum=3.0, value=1.0, label="Embedding Scale"),
#    ],
#    outputs=[gr.Audio(label="Synthesized Audio"), gr.Textbox(label="Message")],
#    title="StyleTTS2 Voice Cloning",
#    description="Upload a reference wav audio file and enter text to synthesize speech."
# )

# Make upper part of the interface as block

def generate_subtitles(audio_file, language):
    # Clear the memory
    gc.collect()

    # Delete temp files if they exist
    if os.path.exists("temp.srt"):
        os.remove("temp.srt")

    output_dir = "subtitles"

    print("Generating subtitles...")
    # Convert the audio file to WAV if necessary
    audio_file = convert_to_wav(audio_file)

    # Define the command for generating subtitles
    whisper_command = f"whisper {audio_file} --model small --language {language} --max_line_count=1 --max_line_width=70 --word_timestamps=True --output_dir ./ --output_format=srt"

    # Run the command
    os.system(whisper_command)
    print("Subtitles generated")

    # Take subtitle file name as after last "/" in audio file name
    subtitle_file = audio_file.split("/")[-1].replace(".wav", ".srt")

    # Correct known mistakes in the subtitle file
    correct_known_mistakes(file_path=subtitle_file)

    return subtitle_file

def prompt_to_video(audio_file, prompt, speed, alpha, beta, diffusion_steps, embedding_scale):
    """
    Handles the entire process of generating a video from a prompt using the existing parse_generate function.
    """
    print("Generating video from prompt...")
    try:
        print("Generating PowerPoint from prompt...")
        pptx_path = slide_generator(prompt)  # Generate PowerPoint from prompt
        print("PowerPoint generated")
        if pptx_path.endswith(".pptx"):
            print("Generating video from PowerPoint...")
            # Use parse_generate with the PowerPoint file
            return parse_generate(
                audio_file=audio_file,  # No reference audio required
                text_input_type="Prompt",
                text_input=None,
                text_file=None,
                srt_input=None,
                zip_file=None,
                pptx_inp=pptx_path,
                speed=speed,
                alpha=alpha,
                beta=beta,
                diffusion_steps=diffusion_steps,
                embedding_scale=embedding_scale
            )
        else:
            return None, None, f"Error: slide_generator did not return a valid PPTX path."
    except Exception as e:
        return None, None, f"Error generating video: {str(e)}"


with gr.Blocks() as iface:
    with gr.Tabs():
        with gr.TabItem("Prompt to Video"):
            with gr.Row():
                with gr.Column():
                    # Input for the reference audio
                    audio_file = gr.Audio(label="Upload a reference audio file", type="filepath")

                    # Input for the prompt
                    prompt_input = gr.Textbox(
                        lines=4,
                        placeholder="Enter a detailed prompt to generate slides",
                        label="Prompt"
                    )
                    # Sliders for customization
                    speed_slider = gr.Slider(0.5, 1.5, value=1.05, label="Speed")
                    alpha_slider = gr.Slider(0.0, 1.0, value=0.3, label="Alpha")
                    beta_slider = gr.Slider(0.0, 1.0, value=0.7, label="Beta")
                    diffusion_slider = gr.Slider(1, 10, value=5, label="Diffusion Steps")
                    embedding_slider = gr.Slider(0.0, 3.0, value=1.0, label="Embedding Scale")

                    # Submit button
                    generate_zip_button = gr.Button(value="Generate Video")

                with gr.Column():
                    # Outputs
                    zip_output = gr.File(label="Download Generated ZIP")
                    message_output = gr.Textbox(label="Message", interactive=False)

            # Connect the button to the new function
            generate_zip_button.click(
                prompt_to_video,
                inputs=[
                    audio_file,  # Reference audio
                    prompt_input,  # Prompt
                    speed_slider,  # Speed
                    alpha_slider,  # Alpha
                    beta_slider,  # Beta
                    diffusion_slider,  # Diffusion steps
                    embedding_slider,  # Embedding scale
                ],
                outputs=[zip_output, message_output]  # Output ZIP file and message
            )
        with gr.TabItem("Voice Cloning"):
            with gr.Row():
                # define inputs
                with gr.Column():
                    audio_file = gr.Audio(label="Upload a reference audio file", type="filepath")
                    txt_radio = gr.Radio(["Plain Text", "TXT File", "SRT File", "ZIP File", "PowerPoint File"],
                                         label="Input Type", value="Plain Text")
                    txt_box = gr.Textbox(lines=2, placeholder="Enter text to synthesize", label="Text to Synthesize",
                                         visible=True)
                    txt_inp = gr.File(label="Upload a TXT file", type="filepath", file_types=[".txt"], visible=False)
                    srt_inp = gr.File(label="Upload an SRT file to generate from subtitles", type="filepath",
                                      file_types=[".srt"], visible=False)
                    zip_inp = gr.File(label="Upload a ZIP file", type="filepath", file_types=[".zip"], visible=False)
                    pptx_inp = gr.File(label="Upload a PowerPoint file to generate from notes", type="filepath",
                                       file_types=[".pptx"], visible=False)

                    with gr.Row():
                        with gr.Column():
                            clear_button = gr.ClearButton(value="Clear")
                        with gr.Column():
                            submit_button = gr.Button(value="Submit")

                    speed_slider = gr.Slider(minimum=0.5, maximum=1.5, value=1.05, label="Speed")
                    alpha_slider = gr.Slider(minimum=0.0, maximum=1.0, value=0.3, label="Alpha")
                    beta_slider = gr.Slider(minimum=0.0, maximum=1.0, value=0.7, label="Beta")
                    diffusion_slider = gr.Slider(minimum=1, maximum=10, value=5, label="Diffusion Steps")
                    embedding_slider = gr.Slider(minimum=0.0, maximum=3.0, value=1.0, label="Embedding Scale")

                with gr.Column():
                    # define outputs
                    audio_output = gr.Audio(label="Synthesized Audio", visible=True)
                    zip_output = gr.File(label="Download ZIP file", visible=False)
                    message_output = gr.Textbox(label="Message")

            txt_radio.change(update_input_fields, inputs=[txt_radio],
                             outputs=[txt_box, txt_inp, srt_inp, zip_inp, pptx_inp, audio_output, zip_output])

            submit_button.click(parse_generate,
                                inputs=[audio_file, txt_radio, txt_box, txt_inp, srt_inp, zip_inp, pptx_inp,
                                        speed_slider, alpha_slider, beta_slider, diffusion_slider, embedding_slider],
                                outputs=[zip_output, audio_output, message_output])

        with gr.TabItem("Subtitle Generator"):
            with gr.Row():
                with gr.Column():
                    subtitle_audio_file = gr.Audio(label="Upload an audio file", type="filepath")
                    language = gr.Dropdown(["English", "Spanish", "French", "German"], label="Language")

                    with gr.Row():
                        with gr.Column():
                            generate_subtitles_button = gr.Button(value="Generate Subtitles")

                with gr.Column():
                    subtitles_output = gr.File(label="Generated Subtitles")

            generate_subtitles_button.click(generate_subtitles, inputs=[subtitle_audio_file, language],
                                            outputs=[subtitles_output])




if __name__ == "__main__":
    iface.queue().launch(server_port=7861, share=True)  # server_name="0.0.0.0",
