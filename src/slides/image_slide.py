from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.slide import Slide

from typing import List, Callable, Optional
from PIL import Image
import random
import tqdm
import os

from src.font import Font
from .slide_utils import add_paragraph

def generate_text_title_image_right(
    presentation: Presentation,
    title: str,
    text: str,
    picture_path: str,
    font:Font, 
    text_font_coeff:float=0.6,
) -> Slide:
    """
    Add a slide with title, text on the left, and picture on the right.

    Args:
    presentation (Presentation): PowerPoint presentation object
    title (str): Title for the slide
    text (str): Text content for the left side of the slide
    picture_path (str): Path to the picture to be inserted on the right side
    font (Font): Font object to manage font styles and paths.
    text_font_coeff (float): Coefficient to adjust the font size of the text relative to the title (default is 0.6).
    Returns:
    Slide
    """

    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)

    slide_height = 9
    slide_width = 16
    margin = min(slide_height, slide_width) / 18

    # image params
    # original image size
    x_pixels, y_pixels = Image.open(picture_path).size
    assert x_pixels == y_pixels or x_pixels < y_pixels, \
        'only vertical and square images can be used'
    # we need image height to be equal to slide height
    image_height = slide_height
    # x_pixels / y_pixels = image_width / image_height
    image_width = x_pixels / y_pixels * image_height
    image_left = slide_width - image_width
    image_top = 0

    slide.shapes.add_picture(
        picture_path,
        left=Inches(image_left),
        top=Inches(image_top),
        width=Inches(image_width),
        height=Inches(image_height),
    )

    # title params
    title_left = margin
    title_top = margin
    title_width = slide_width - image_width - 2*margin
    title_height = slide_height / 6

    title_box = slide.shapes.add_textbox(
        left=Inches(title_left),
        top=Inches(title_top),
        width=Inches(title_width),
        height=Inches(title_height),
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.word_wrap = False
    
    # title_paragraph = title_frame.add_paragraph()
    title_paragraph = add_paragraph(title_frame)
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.text = title

    for max_size in range(font.max_size)[::-5]: 
        try: 
            title_frame.fit_text(
                font_file=font.bold, 
                max_size=max_size,
                bold=True,
            )
            break
        except: 
            pass

    # text params
    title_left = margin
    text_top = title_height + margin*2
    text_width = slide_width - image_width - 2*margin
    text_height = slide_height - title_height - 3*margin

    text_box = slide.shapes.add_textbox(
        left=Inches(title_left),
        top=Inches(text_top),
        width=Inches(text_width),
        height=Inches(text_height),
    )
    text_frame = text_box.text_frame
    text_frame.clear()
    
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = False
    
    # text_paragraph = text_frame.add_paragraph()
    text_paragraph = add_paragraph(text_frame)
    text_paragraph.text = text 
    text_paragraph.alignment = PP_ALIGN.CENTER

    for max_size in range(int(max_size*text_font_coeff))[::-5]:
        try: 
            text_frame.fit_text(font_file=font.basic, max_size=max_size)
            break
        except: 
            pass
            
    
    return slide

def generate_text_title_image_left(
    presentation: Presentation,
    title: str,
    text: str,
    picture_path: str,
    font:Font,
    text_font_coeff:float=0.6,
) -> Slide:
    """
    Add a slide with title, text on the left, and picture on the right.

    Args:
        presentation (Presentation): PowerPoint presentation object
        title (str): Title for the slide
        text (str): Text content for the left side of the slide
        picture_path (str): Path to the picture to be inserted on the right side
        font (Font): Font object to manage font styles and paths.
        text_font_coeff (float): Coefficient to adjust the font 
            size of the text relative to the title (default is 0.6).
    
    Returns:
        Slide
    """

    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)

    slide_height = 9
    slide_width = 16
    margin = min(slide_height, slide_width) / 18

    # image params
    # original image size
    x_pixels, y_pixels = Image.open(picture_path).size
    assert x_pixels == y_pixels or x_pixels < y_pixels, \
        'only vertical and square images can be used'
    # we need image height to be equal to slide height
    image_height = slide_height
    # x_pixels / y_pixels = image_width / image_height
    image_width = x_pixels / y_pixels * image_height
    image_left = 0
    image_top = 0

    slide.shapes.add_picture(
        picture_path,
        left=Inches(image_left),
        top=Inches(image_top),
        width=Inches(image_width),
        height=Inches(image_height),
    )
    
    # title params
    title_left = image_width + margin 
    title_top = margin
    title_width = slide_width - image_width - 2 * margin
    title_height = slide_height / 6

    title_box = slide.shapes.add_textbox(
        left=Inches(title_left),
        top=Inches(title_top),
        width=Inches(title_width),
        height=Inches(title_height),
    )
    title_frame = title_box.text_frame
    title_frame.clear()

    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.word_wrap = False
    
    # title_paragraph = title_frame.add_paragraph()
    title_paragraph = add_paragraph(title_frame)
    title_paragraph.text = title
    title_paragraph.alignment = PP_ALIGN.CENTER

    for max_size in range(font.max_size)[::-5]: 
        try: 
            title_frame.fit_text(
                font_file=font.bold, 
                max_size=max_size,
                bold=True,
            )
            break
        except: 
            pass
            
    # text params
    text_left = title_left
    text_top = title_height + margin * 2
    text_width = slide_width - image_width - 2 * margin
    text_height = slide_height - title_height - 3 * margin

    text_box = slide.shapes.add_textbox(
        left=Inches(text_left),
        top=Inches(text_top),
        width=Inches(text_width),
        height=Inches(text_height),
    )
    text_frame = text_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = False
    
    # text_paragraph = text_frame.add_paragraph()
    text_paragraph = add_paragraph(text_frame)
    text_paragraph.text = text
    text_paragraph.alignment = PP_ALIGN.CENTER
    
    for max_size in range(int(max_size*text_font_coeff))[::-5]:
        try: 
            text_frame.fit_text(
                font_file=font.basic, 
                max_size=max_size
            )
            break
        except: 
            pass
           
    return slide

def generate_image_slide(
    presentation: Presentation, 
    title: str, 
    text: str, 
    picture_path: str,
    font: Font,
    text_font_coeff: float = 0.6,
) -> Slide:
    """
    Generate a slide with a title, text, and an image.

    This function creates a slide in a PowerPoint presentation that includes a title, 
    text, and an image. The layout is chosen randomly between two options: 
    image on the right or image on the left.

    Args:
        presentation (Presentation): PowerPoint presentation object.
        title (str): Title for the slide.
        text (str): Text content for the slide.
        picture_path (str): Path to the picture to be inserted in the slide.
        font (Font): Font object to manage font styles and paths.
        text_font_coeff (float, optional): Coefficient to adjust the font size of the text 
                                           relative to the title (default is 0.65).

    Returns:
        Slide
    """
    gen_func = random.choice([
        generate_text_title_image_right,
        generate_text_title_image_left,
    ])
    return gen_func(
        presentation=presentation,
        title=title,
        text=text,
        picture_path=picture_path,
        font=font,
        text_font_coeff=text_font_coeff,
    )