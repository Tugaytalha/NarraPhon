from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

import random
import os
from PIL import Image
from typing import List, Callable

import tqdm

def add_paragraph(text_frame): 
    try:
        title_paragraph = text_frame.paragraphs[0]
    except:
        title_paragraph = text_frame.add_paragraph()
    return title_paragraph
    
def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def set_shape_transparency(shape, alpha):
    """ Set the transparency (alpha) of a shape"""
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    SubElement(sF, 'a:alpha', val=str(int(alpha*100000)))