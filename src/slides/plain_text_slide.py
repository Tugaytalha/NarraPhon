from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.slide import Slide

from typing import List, Callable, Optional
from PIL import Image
import random
import tqdm
import os

from src.font import Font
from .slide_utils import add_paragraph

def generate_plain_text_slide(
    presentation: Presentation,
    title: str,
    text: str,
    background_path: str,
    font: Font,
    text_font_coeff: float = 0.6,
) -> Slide:
    """
    Add a slide with title and text on a background image.

    Args:
        presentation (Presentation): PowerPoint presentation object
        title (str): Title for the slide
        text (str): Text content for the slide
        background_path (str): Path to the background image
        font (Font): Font object to manage font styles and paths
        text_font_coeff (float): Coefficient to adjust text font size relative to title

    Returns:
        Slide: The created slide object
    """
    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)

    slide_height = 9
    slide_width = 16
    margin = min(slide_height, slide_width) / 18

    # Add background image
    if background_path:
        pic = slide.shapes.add_picture(
            background_path, 0, 0,
            width=presentation.slide_width,
            height=presentation.slide_height
        )
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

    # Add title
    title_left = margin
    title_top = margin
    title_width = slide_width - 2 * margin
    title_height = slide_height / 6

    title_box = slide.shapes.add_textbox(
        left=Inches(title_left),
        top=Inches(title_top),
        width=Inches(title_width),
        height=Inches(title_height)
    )
    
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_frame.word_wrap = True
    
    title_paragraph = add_paragraph(title_frame)
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.text = title

    for max_size in range(font.max_size)[::-5]:
        try:
            title_frame.fit_text(
                font_file=font.bold,
                max_size=max_size,
                bold=True
            )
            break
        except:
            pass

    # Add text
    text_left = margin
    text_top = title_height + margin * 2
    text_width = slide_width - 2 * margin
    text_height = slide_height - title_height - 3 * margin

    text_box = slide.shapes.add_textbox(
        left=Inches(text_left),
        top=Inches(text_top),
        width=Inches(text_width),
        height=Inches(text_height)
    )
    
    text_frame = text_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.word_wrap = True
    
    text_paragraph = add_paragraph(text_frame)
    text_paragraph.alignment = PP_ALIGN.CENTER
    text_paragraph.text = text

    for max_size in range(int(max_size * text_font_coeff))[::-5]:
        try:
            text_frame.fit_text(
                font_file=font.basic,
                max_size=max_size
            )
            break
        except:
            pass

    return slide
    