from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.slide import Slide

from typing import List, Callable, Optional, Tuple
from PIL import Image
import random
import tqdm
import os


from .image_slide import generate_image_slide
from .plain_text_slide import generate_plain_text_slide
from .title_slide import generate_title_slide

from src.font import Font

def generate_slide(
    presentation: Presentation,
    title: str,
    text: Optional[Tuple[str, str]] = None,
    background_path: Optional[str] = None,
    picture_path: Optional[str] = None,
    font: Font = None, 
    text_font_coeff: float = 0.6,
) -> None:
    """
    Generate a slide in the presentation based on the provided content.
    
    Args:
        text (Optional[Tuple[str, str]]): Tuple of (slide_text, speaker_notes)
    """
    slide_text = None if text is None else text[0]
    speaker_notes = None if text is None else text[1]
    
    # All slides will use image_slide layout
    slide = generate_image_slide(
        presentation=presentation,
        title=title,
        text=slide_text,
        picture_path=picture_path,
        font=font,
        text_font_coeff=text_font_coeff,
    )

    # Add speaker notes if they exist
    if speaker_notes and slide is not None:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = speaker_notes

    return slide