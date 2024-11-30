from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.slide import Slide

import random
import os
from PIL import Image
from typing import List, Callable

import tqdm

from .slide_utils import set_shape_transparency, add_paragraph

from src.font import Font


def generate_title_slide(
    presentation: Presentation,
    title: str,
    font: Font, 
    background_path: str = None, 
) -> Slide:
    """
    Add a slide with title, text placeholders on the blurred background image.

    Args:
        presentation (Presentation): PowerPoint presentation object
        title (str): Title for the slide
        text (str): Text content for the slide
        background_path (str): Path to the background image for the slide
        font (Font): Font object to manage font styles and paths.
    Returns:
        Slide: The created slide object
    """

    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)

    slide_height = 9
    slide_width = 16
    margin = min(slide_height, slide_width) / 18

    # Background image
    if background_path:
        pic = slide.shapes.add_picture(
            background_path, 0, 0,
            width=presentation.slide_width,
            height=presentation.slide_height
        )
        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

    # Title placeholder
    title_left = margin
    title_top = margin
    title_width = slide_width - 2 * margin
    title_height = slide_height - 2 * margin

    title_box = slide.shapes.add_textbox(
        left=Inches(title_left),
        top=Inches(title_top),
        width=Inches(title_width),
        height=Inches(title_height),
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    title_frame.word_wrap = False

    title_paragraph = add_paragraph(title_frame)
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.text = title
    
    for max_size in range(font.max_size, 0, -5): 
        try: 
            title_frame.fit_text(
                font_file=font.bold, 
                max_size=max_size,
                bold=True,
            )
            break
        except TypeError: 
            pass

    # Set white color and transparency to title shape
    title_fill = title_box.fill
    title_fill.solid()
    title_fill.fore_color.rgb = RGBColor(255, 255, 255)
    set_shape_transparency(title_box, 0.5)

    return slide